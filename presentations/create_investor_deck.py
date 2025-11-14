#!/usr/bin/env python3
"""
Create a beautiful, professional investor pitch deck for Hospital Management System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
PRIMARY_BLUE = RGBColor(30, 58, 138)  # #1E3A8A - Dark blue
ACCENT_BLUE = RGBColor(59, 130, 246)  # #3B82F6 - Bright blue
TEAL = RGBColor(13, 148, 136)  # #0D9488
GREEN = RGBColor(16, 185, 129)  # #10B981
ORANGE = RGBColor(249, 115, 22)  # #F97316
RED = RGBColor(239, 68, 68)  # #EF4444
DARK_GRAY = RGBColor(31, 41, 55)  # #1F2937
LIGHT_GRAY = RGBColor(243, 244, 246)  # #F3F4F6
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add background rectangle
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BLUE
    bg.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4),
        Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = LIGHT_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title):
    """Add a content slide with header"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Add title in header
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    return slide

def add_bullet_box(slide, left, top, width, height, bullets, font_size=16):
    """Add a text box with bullets"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_after = Pt(12)

    return text_box

def add_section_slide(prs, section_title):
    """Add a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    return slide

def create_investor_deck():
    """Create the complete investor pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Revolutionizing Hospital Operations",
        "Real-Time Intelligent Hospital Management System"
    )

    # Slide 2: The Problem
    slide = add_content_slide(prs, "The Problem: Hospitals Face Critical Inefficiencies")

    # Problem statistics boxes
    problems = [
        ("60-90 min", "Average ED Wait Time", RED),
        ("$265B", "Wasted Annually in US", ORANGE),
        ("25-30%", "Staff Time on Paperwork", ORANGE),
        ("Hours", "Critical Result Delays", RED)
    ]

    x_positions = [Inches(0.5), Inches(2.8), Inches(5.1), Inches(7.4)]
    for i, (stat, desc, color) in enumerate(problems):
        # Stat box
        stat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], Inches(1.5),
            Inches(2.2), Inches(1.8)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = color
        stat_box.line.fill.background()

        # Stat text
        stat_frame = stat_box.text_frame
        stat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = stat_frame.paragraphs[0]
        p.text = stat
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description below
        desc_box = slide.shapes.add_textbox(
            x_positions[i], Inches(3.4),
            Inches(2.2), Inches(0.6)
        )
        desc_frame = desc_box.text_frame
        dp = desc_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = DARK_GRAY
        dp.alignment = PP_ALIGN.CENTER

    # Real-world impact section
    impact_bullets = [
        "Delayed critical test results lead to adverse patient outcomes",
        "Manual billing processes result in 15-20% revenue leakage",
        "Poor queue management causes patient walkouts and lost revenue",
        "Disconnected systems require duplicate data entry, increasing errors"
    ]
    add_bullet_box(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(2.5), impact_bullets, 14)

    # Slide 3: The Solution
    slide = add_content_slide(prs, "The Solution: Real-Time Hospital Management Platform")

    # Value proposition box
    value_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(1.2)
    )
    value_box.fill.solid()
    value_box.fill.fore_color.rgb = LIGHT_GRAY
    value_box.line.fill.background()

    value_text = value_box.text_frame
    vp = value_text.paragraphs[0]
    vp.text = "Enterprise-grade system that optimizes every aspect of hospital operations through real-time workflow automation"
    vp.font.size = Pt(18)
    vp.font.color.rgb = PRIMARY_BLUE
    vp.font.bold = True
    vp.alignment = PP_ALIGN.CENTER
    value_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Key differentiators
    differentiators = [
        ("Real-Time Engine", "WebSocket-based instant notifications (<50ms) vs. delayed batch updates"),
        ("Intelligent Automation", "AI-driven workflow optimization vs. manual processes"),
        ("Unified Platform", "Single integrated system vs. fragmented point solutions"),
        ("Modern Architecture", "Cloud-native, scalable vs. legacy on-premise systems")
    ]

    y_pos = Inches(3)
    for title, desc in differentiators:
        # Bullet point
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7), y_pos + Inches(0.05),
            Inches(0.15), Inches(0.15)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = TEAL
        bullet.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), y_pos,
            Inches(8.5), Inches(0.3)
        )
        t_frame = title_box.text_frame
        tp = t_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = PRIMARY_BLUE

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(1), y_pos + Inches(0.3),
            Inches(8.5), Inches(0.4)
        )
        d_frame = desc_box.text_frame
        dp = d_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(13)
        dp.font.color.rgb = DARK_GRAY

        y_pos += Inches(0.9)

    # Slide 4: Product Overview
    slide = add_content_slide(prs, "Product Overview: 8 Integrated Modules")

    modules = [
        ("Reception", "Patient registration, check-in, queue management", "40% wait time reduction"),
        ("Doctor", "EMR, consultations, prescriptions, test ordering", "2 hrs/day saved"),
        ("Laboratory", "Test management, critical alerts, results", "Instant critical delivery"),
        ("Pharmacy", "Prescription queue, dispensing, inventory", "95% fill rate"),
        ("Radiology", "Imaging orders, scheduling, reporting", "Streamlined workflow"),
        ("Billing", "Automated billing, insurance, payments", "15% revenue recovery"),
        ("Analytics", "Real-time dashboards, KPIs, reports", "Data-driven decisions"),
        ("Notifications", "Role-based alerts, critical notifications", "Zero missed results")
    ]

    # Create 2 columns of 4 modules each
    for i, (module, features, value) in enumerate(modules):
        col = i // 4
        row = i % 4

        x = Inches(0.5) if col == 0 else Inches(5.3)
        y = Inches(1.5) + Inches(1.45) * row

        # Module box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(4.5), Inches(1.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(2)

        # Module name
        name_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.15),
            Inches(4.1), Inches(0.3)
        )
        n_frame = name_box.text_frame
        np = n_frame.paragraphs[0]
        np.text = module
        np.font.size = Pt(16)
        np.font.bold = True
        np.font.color.rgb = PRIMARY_BLUE

        # Features
        feat_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.5),
            Inches(4.1), Inches(0.4)
        )
        f_frame = feat_box.text_frame
        fp = f_frame.paragraphs[0]
        fp.text = features
        fp.font.size = Pt(11)
        fp.font.color.rgb = DARK_GRAY

        # Value tag
        val_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.2), y + Inches(0.95),
            Inches(4.1), Inches(0.25)
        )
        val_box.fill.solid()
        val_box.fill.fore_color.rgb = GREEN
        val_box.line.fill.background()

        v_frame = val_box.text_frame
        vp = v_frame.paragraphs[0]
        vp.text = f"✓ {value}"
        vp.font.size = Pt(10)
        vp.font.bold = True
        vp.font.color.rgb = WHITE
        vp.alignment = PP_ALIGN.CENTER
        v_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Slide 5: Market Opportunity
    slide = add_content_slide(prs, "Market Opportunity: $68B Healthcare IT Market")

    # Big number
    market_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(4), Inches(1.5)
    )
    m_frame = market_box.text_frame
    mp = m_frame.paragraphs[0]
    mp.text = "$68B"
    mp.font.size = Pt(72)
    mp.font.bold = True
    mp.font.color.rgb = ACCENT_BLUE
    mp.alignment = PP_ALIGN.CENTER
    m_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Growth rate
    growth_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3),
        Inches(4), Inches(0.5)
    )
    g_frame = growth_box.text_frame
    gp = g_frame.paragraphs[0]
    gp.text = "Growing at 13.7% CAGR"
    gp.font.size = Pt(20)
    gp.font.color.rgb = GREEN
    gp.font.bold = True
    gp.alignment = PP_ALIGN.CENTER

    # Market breakdown
    markets = [
        ("TAM", "Total Addressable Market", "$38B", "Hospital Information Systems Global"),
        ("SAM", "Serviceable Addressable Market", "$375M - $1.25B", "Mid-to-large US hospitals (200+ beds)"),
        ("SOM", "Serviceable Obtainable Market", "$7.5M - $25M", "2% market share in 3 years (50 hospitals)")
    ]

    y_pos = Inches(1.5)
    for abbr, name, value, desc in markets:
        # Abbreviation badge
        abbr_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5), y_pos,
            Inches(0.8), Inches(0.8)
        )
        abbr_box.fill.solid()
        abbr_box.fill.fore_color.rgb = PRIMARY_BLUE
        abbr_box.line.fill.background()

        a_frame = abbr_box.text_frame
        ap = a_frame.paragraphs[0]
        ap.text = abbr
        ap.font.size = Pt(18)
        ap.font.bold = True
        ap.font.color.rgb = WHITE
        ap.alignment = PP_ALIGN.CENTER
        a_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Name and value
        text_box = slide.shapes.add_textbox(
            Inches(6), y_pos,
            Inches(3.5), Inches(0.8)
        )
        t_frame = text_box.text_frame

        # Name
        tp = t_frame.paragraphs[0]
        tp.text = name
        tp.font.size = Pt(12)
        tp.font.color.rgb = DARK_GRAY

        # Value
        vp = t_frame.add_paragraph()
        vp.text = value
        vp.font.size = Pt(18)
        vp.font.bold = True
        vp.font.color.rgb = ACCENT_BLUE

        # Description
        dp = t_frame.add_paragraph()
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = DARK_GRAY

        y_pos += Inches(1.7)

    # Market drivers
    drivers_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(9), Inches(0.4)
    )
    dt_frame = drivers_title.text_frame
    dtp = dt_frame.paragraphs[0]
    dtp.text = "Market Drivers:"
    dtp.font.size = Pt(16)
    dtp.font.bold = True
    dtp.font.color.rgb = PRIMARY_BLUE

    drivers = [
        "Mandated EHR adoption (95% of hospitals now digital)",
        "Post-pandemic focus on operational efficiency",
        "Value-based care reimbursement models",
        "Patient experience as competitive differentiator"
    ]

    add_bullet_box(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(2.5), drivers, 12)

    # Slide 6: Business Model
    slide = add_content_slide(prs, "Business Model: High-Margin SaaS Revenue")

    # Pricing tiers
    tiers = [
        ("Essential", "50-150 beds", "$50K/year", "Core modules, 100 users", TEAL),
        ("Professional", "150-300 beds", "$150K/year", "All modules, 300 users, analytics", ACCENT_BLUE),
        ("Enterprise", "300+ beds", "$500K/year", "Unlimited users, custom integrations", PRIMARY_BLUE)
    ]

    x_start = Inches(0.5)
    for i, (tier, size, price, features, color) in enumerate(tiers):
        x = x_start + Inches(3.1) * i

        # Tier box
        tier_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.5),
            Inches(2.9), Inches(2.5)
        )
        tier_box.fill.solid()
        tier_box.fill.fore_color.rgb = color
        tier_box.line.fill.background()

        # Content
        content_box = slide.shapes.add_textbox(
            x + Inches(0.2), Inches(1.7),
            Inches(2.5), Inches(2.1)
        )
        c_frame = content_box.text_frame

        # Tier name
        cp1 = c_frame.paragraphs[0]
        cp1.text = tier
        cp1.font.size = Pt(20)
        cp1.font.bold = True
        cp1.font.color.rgb = WHITE
        cp1.alignment = PP_ALIGN.CENTER

        # Size
        cp2 = c_frame.add_paragraph()
        cp2.text = size
        cp2.font.size = Pt(12)
        cp2.font.color.rgb = LIGHT_GRAY
        cp2.alignment = PP_ALIGN.CENTER
        cp2.space_before = Pt(6)

        # Price
        cp3 = c_frame.add_paragraph()
        cp3.text = price
        cp3.font.size = Pt(28)
        cp3.font.bold = True
        cp3.font.color.rgb = WHITE
        cp3.alignment = PP_ALIGN.CENTER
        cp3.space_before = Pt(12)

        # Features
        cp4 = c_frame.add_paragraph()
        cp4.text = features
        cp4.font.size = Pt(11)
        cp4.font.color.rgb = LIGHT_GRAY
        cp4.alignment = PP_ALIGN.CENTER
        cp4.space_before = Pt(12)

    # Unit economics
    metrics_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.3),
        Inches(9), Inches(0.4)
    )
    mt_frame = metrics_title.text_frame
    mtp = mt_frame.paragraphs[0]
    mtp.text = "Unit Economics:"
    mtp.font.size = Pt(18)
    mtp.font.bold = True
    mtp.font.color.rgb = PRIMARY_BLUE

    metrics = [
        ("LTV/CAC Ratio", "18-100x", GREEN),
        ("Gross Margin", "75-85%", GREEN),
        ("Payback Period", "3-6 months", GREEN)
    ]

    x_pos = Inches(0.5)
    for metric, value, color in metrics:
        # Metric box
        m_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, Inches(4.9),
            Inches(2.9), Inches(1.2)
        )
        m_box.fill.solid()
        m_box.fill.fore_color.rgb = LIGHT_GRAY
        m_box.line.color.rgb = color
        m_box.line.width = Pt(3)

        # Content
        m_content = slide.shapes.add_textbox(
            x_pos + Inches(0.2), Inches(5),
            Inches(2.5), Inches(1)
        )
        mc_frame = m_content.text_frame
        mc_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Metric name
        mcp1 = mc_frame.paragraphs[0]
        mcp1.text = metric
        mcp1.font.size = Pt(12)
        mcp1.font.color.rgb = DARK_GRAY
        mcp1.alignment = PP_ALIGN.CENTER

        # Value
        mcp2 = mc_frame.add_paragraph()
        mcp2.text = value
        mcp2.font.size = Pt(24)
        mcp2.font.bold = True
        mcp2.font.color.rgb = color
        mcp2.alignment = PP_ALIGN.CENTER

        x_pos += Inches(3.1)

    # Slide 7: Financial Projections
    slide = add_content_slide(prs, "Financial Projections: Path to $25M ARR in 3 Years")

    # Revenue table
    years_data = [
        ("Year 1", "10", "10", "$1.5M", "-"),
        ("Year 2", "25", "35", "$7M", "367%"),
        ("Year 3", "40", "75", "$18M", "157%"),
        ("Year 4", "60", "135", "$35M", "94%"),
        ("Year 5", "80", "215", "$60M", "71%")
    ]

    # Table header
    headers = ["Year", "New", "Total", "ARR", "Growth"]
    x_positions = [Inches(0.5), Inches(2.5), Inches(4), Inches(5.8), Inches(7.8)]

    # Header row
    for i, header in enumerate(headers):
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], Inches(1.5),
            Inches(1.6), Inches(0.4)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = PRIMARY_BLUE
        header_box.line.fill.background()

        h_frame = header_box.text_frame
        h_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        hp = h_frame.paragraphs[0]
        hp.text = header
        hp.font.size = Pt(12)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER

    # Data rows
    y_pos = Inches(2)
    for year, new, total, arr, growth in years_data:
        row_data = [year, new, total, arr, growth]

        for i, value in enumerate(row_data):
            data_box = slide.shapes.add_textbox(
                x_positions[i], y_pos,
                Inches(1.6), Inches(0.5)
            )
            d_frame = data_box.text_frame
            d_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            dp = d_frame.paragraphs[0]
            dp.text = value
            dp.font.size = Pt(14)
            dp.font.color.rgb = DARK_GRAY
            dp.alignment = PP_ALIGN.CENTER

            # Highlight ARR column
            if i == 3:
                dp.font.bold = True
                dp.font.color.rgb = ACCENT_BLUE
                dp.font.size = Pt(16)

        y_pos += Inches(0.6)

    # Key assumptions
    assumptions_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5),
        Inches(4.5), Inches(0.3)
    )
    at_frame = assumptions_title.text_frame
    atp = at_frame.paragraphs[0]
    atp.text = "Key Assumptions:"
    atp.font.size = Pt(14)
    atp.font.bold = True
    atp.font.color.rgb = PRIMARY_BLUE

    assumptions = [
        "Avg contract: $150K/year",
        "90% retention rate",
        "20% annual price increases"
    ]

    y = Inches(5.9)
    for assumption in assumptions:
        a_box = slide.shapes.add_textbox(
            Inches(0.7), y,
            Inches(4), Inches(0.3)
        )
        a_frame = a_box.text_frame
        ap = a_frame.paragraphs[0]
        ap.text = f"• {assumption}"
        ap.font.size = Pt(11)
        ap.font.color.rgb = DARK_GRAY
        y += Inches(0.3)

    # Break-even callout
    breakeven_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(5.5),
        Inches(4), Inches(1.3)
    )
    breakeven_box.fill.solid()
    breakeven_box.fill.fore_color.rgb = GREEN
    breakeven_box.line.fill.background()

    be_content = slide.shapes.add_textbox(
        Inches(5.7), Inches(5.7),
        Inches(3.6), Inches(0.9)
    )
    be_frame = be_content.text_frame
    be_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    be_p1 = be_frame.paragraphs[0]
    be_p1.text = "Break-Even"
    be_p1.font.size = Pt(16)
    be_p1.font.bold = True
    be_p1.font.color.rgb = WHITE
    be_p1.alignment = PP_ALIGN.CENTER

    be_p2 = be_frame.add_paragraph()
    be_p2.text = "Month 18 at 25 customers"
    be_p2.font.size = Pt(14)
    be_p2.font.color.rgb = LIGHT_GRAY
    be_p2.alignment = PP_ALIGN.CENTER

    # Slide 8: The Ask
    slide = add_content_slide(prs, "The Ask: $2M Seed Round")

    # Main ask box
    ask_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(8), Inches(2)
    )
    ask_box.fill.solid()
    ask_box.fill.fore_color.rgb = PRIMARY_BLUE
    ask_box.line.fill.background()

    ask_content = slide.shapes.add_textbox(
        Inches(1.3), Inches(1.8),
        Inches(7.4), Inches(1.4)
    )
    ac_frame = ask_content.text_frame

    # Amount
    ac_p1 = ac_frame.paragraphs[0]
    ac_p1.text = "$2 Million Seed Round"
    ac_p1.font.size = Pt(40)
    ac_p1.font.bold = True
    ac_p1.font.color.rgb = WHITE
    ac_p1.alignment = PP_ALIGN.CENTER

    # Valuation
    ac_p2 = ac_frame.add_paragraph()
    ac_p2.text = "$8M post-money valuation"
    ac_p2.font.size = Pt(20)
    ac_p2.font.color.rgb = LIGHT_GRAY
    ac_p2.alignment = PP_ALIGN.CENTER
    ac_p2.space_before = Pt(12)

    # Use of funds
    allocation = [
        ("Product Development", "30%", "$600K", "3 engineers, UI/UX, QA", ACCENT_BLUE),
        ("Sales & Marketing", "35%", "$700K", "2 sales reps, marketing, conferences", TEAL),
        ("Customer Success", "15%", "$300K", "2 implementation specialists", GREEN),
        ("Infrastructure", "10%", "$200K", "AWS hosting, tools, security", ORANGE),
        ("Operations", "10%", "$200K", "Office, legal, accounting", DARK_GRAY)
    ]

    y_pos = Inches(3.8)
    for category, pct, amount, desc, color in allocation:
        # Category box
        cat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y_pos,
            Inches(9), Inches(0.5)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = LIGHT_GRAY
        cat_box.line.color.rgb = color
        cat_box.line.width = Pt(2)

        # Category name
        name_box = slide.shapes.add_textbox(
            Inches(0.7), y_pos + Inches(0.05),
            Inches(3), Inches(0.4)
        )
        n_frame = name_box.text_frame
        n_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = n_frame.paragraphs[0]
        np.text = category
        np.font.size = Pt(13)
        np.font.bold = True
        np.font.color.rgb = color

        # Percentage
        pct_box = slide.shapes.add_textbox(
            Inches(4), y_pos + Inches(0.05),
            Inches(0.8), Inches(0.4)
        )
        p_frame = pct_box.text_frame
        p_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = p_frame.paragraphs[0]
        pp.text = pct
        pp.font.size = Pt(13)
        pp.font.bold = True
        pp.font.color.rgb = color
        pp.alignment = PP_ALIGN.RIGHT

        # Amount
        amt_box = slide.shapes.add_textbox(
            Inches(5), y_pos + Inches(0.05),
            Inches(1), Inches(0.4)
        )
        a_frame = amt_box.text_frame
        a_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        ap = a_frame.paragraphs[0]
        ap.text = amount
        ap.font.size = Pt(13)
        ap.font.bold = True
        ap.font.color.rgb = DARK_GRAY
        ap.alignment = PP_ALIGN.RIGHT

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(6.2), y_pos + Inches(0.05),
            Inches(3), Inches(0.4)
        )
        d_frame = desc_box.text_frame
        d_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        dp = d_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = DARK_GRAY

        y_pos += Inches(0.6)

    # Slide 9: Investment Highlights
    slide = add_content_slide(prs, "Why Invest Now?")

    highlights = [
        ("1", "Large Growing Market", "$68B healthcare IT market at 13.7% CAGR", ACCENT_BLUE),
        ("2", "Proven Product-Market Fit", "Production-ready platform addressing validated pain points", TEAL),
        ("3", "Massive Cost Advantage", "10x cheaper, 20x faster implementation than legacy systems", GREEN),
        ("4", "Recurring Revenue Model", "SaaS with 90%+ retention and 18-100x LTV/CAC", ORANGE),
        ("5", "Scalable Technology", "Cloud-native architecture enabling rapid growth", PRIMARY_BLUE),
        ("6", "Clear Path to Exit", "Acquisition by Epic, Cerner, or IPO with $100M+ ARR", RED)
    ]

    y_pos = Inches(1.5)
    for num, title, desc, color in highlights:
        # Number badge
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7), y_pos,
            Inches(0.5), Inches(0.5)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = color
        num_box.line.fill.background()

        n_frame = num_box.text_frame
        n_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = n_frame.paragraphs[0]
        np.text = num
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.4), y_pos,
            Inches(8), Inches(0.3)
        )
        t_frame = title_box.text_frame
        tp = t_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = color

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(1.4), y_pos + Inches(0.3),
            Inches(8), Inches(0.5)
        )
        d_frame = desc_box.text_frame
        dp = d_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = DARK_GRAY

        y_pos += Inches(0.95)

    # Slide 10: Contact / Thank You
    slide = add_title_slide(prs, "Let's Transform Healthcare Together", "Thank You")

    # Contact info box
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(5.5),
        Inches(6), Inches(1.5)
    )
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = WHITE
    contact_box.line.color.rgb = LIGHT_GRAY
    contact_box.line.width = Pt(2)
    contact_box.shadow.inherit = False

    contact_content = slide.shapes.add_textbox(
        Inches(2.3), Inches(5.7),
        Inches(5.4), Inches(1.1)
    )
    cc_frame = contact_content.text_frame

    # Email
    cc_p1 = cc_frame.paragraphs[0]
    cc_p1.text = "📧 investors@hospital-management-system.com"
    cc_p1.font.size = Pt(16)
    cc_p1.font.color.rgb = PRIMARY_BLUE
    cc_p1.alignment = PP_ALIGN.CENTER

    # Phone
    cc_p2 = cc_frame.add_paragraph()
    cc_p2.text = "📱 +1 (555) 123-4567"
    cc_p2.font.size = Pt(16)
    cc_p2.font.color.rgb = PRIMARY_BLUE
    cc_p2.alignment = PP_ALIGN.CENTER
    cc_p2.space_before = Pt(8)

    # Website
    cc_p3 = cc_frame.add_paragraph()
    cc_p3.text = "🌐 www.hospital-management-system.com"
    cc_p3.font.size = Pt(16)
    cc_p3.font.color.rgb = PRIMARY_BLUE
    cc_p3.alignment = PP_ALIGN.CENTER
    cc_p3.space_before = Pt(8)

    # Save presentation
    prs.save('01-INVESTOR-PITCH-DECK.pptx')
    print("✓ Investor Pitch Deck created successfully!")

if __name__ == "__main__":
    create_investor_deck()
