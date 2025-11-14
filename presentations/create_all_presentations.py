#!/usr/bin/env python3
"""
Create beautiful, professional presentations for all stakeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
PRIMARY_BLUE = RGBColor(30, 58, 138)
ACCENT_BLUE = RGBColor(59, 130, 246)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(16, 185, 129)
ORANGE = RGBColor(249, 115, 22)
RED = RGBColor(239, 68, 68)
DARK_GRAY = RGBColor(31, 41, 55)
LIGHT_GRAY = RGBColor(243, 244, 246)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle, bg_color=PRIMARY_BLUE):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = LIGHT_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bg_color=PRIMARY_BLUE):
    """Add a content slide with header"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = bg_color
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    return slide

def create_hospital_admin_presentation():
    """Create Hospital Administrators Executive Presentation"""
    print("Creating Hospital Administrators presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Transform Your Hospital Operations",
                   "Digital Innovation for Operational Excellence")

    # Slide 2: The Challenge
    slide = add_content_slide(prs, "Healthcare Challenges Today")

    challenges = [
        ("60-90min", "Patient Wait Times", RED),
        ("$1.4M", "Annual Inefficiency Cost", ORANGE),
        ("25%", "Staff Time Wasted", ORANGE),
        ("15-20%", "Revenue Leakage", RED)
    ]

    x_positions = [Inches(0.5), Inches(2.8), Inches(5.1), Inches(7.4)]
    for i, (stat, desc, color) in enumerate(challenges):
        stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         x_positions[i], Inches(1.5), Inches(2.2), Inches(1.8))
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = color
        stat_box.line.fill.background()

        stat_frame = stat_box.text_frame
        stat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = stat_frame.paragraphs[0]
        p.text = stat
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        desc_box = slide.shapes.add_textbox(x_positions[i], Inches(3.4), Inches(2.2), Inches(0.6))
        desc_frame = desc_box.text_frame
        dp = desc_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = DARK_GRAY
        dp.alignment = PP_ALIGN.CENTER

    # Slide 3: ROI Overview
    slide = add_content_slide(prs, "ROI: $1.4M Annual Benefits")

    # Big ROI number
    roi_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(2))
    roi_frame = roi_box.text_frame
    roi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = roi_frame.paragraphs[0]
    p1.text = "$1.4M"
    p1.font.size = Pt(72)
    p1.font.bold = True
    p1.font.color.rgb = GREEN
    p1.alignment = PP_ALIGN.CENTER

    p2 = roi_frame.add_paragraph()
    p2.text = "Annual Benefits"
    p2.font.size = Pt(24)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # Breakdown
    breakdown = [
        ("Revenue Improvements", "$750K/year", GREEN),
        ("Cost Savings", "$650K/year", TEAL)
    ]

    y_pos = Inches(4)
    for label, value, color in breakdown:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1.5), y_pos, Inches(7), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(3)

        text_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.1), Inches(6.6), Inches(0.6))
        t_frame = text_box.text_frame
        t_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        tp1 = t_frame.paragraphs[0]
        tp1.text = f"{label}: "
        tp1.font.size = Pt(18)
        tp1.font.color.rgb = DARK_GRAY

        run = tp1.runs[0]
        run2 = tp1.add_run()
        run2.text = value
        run2.font.size = Pt(24)
        run2.font.bold = True
        run2.font.color.rgb = color

        y_pos += Inches(1)

    # Payback period
    payback_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(3), Inches(6.2), Inches(4), Inches(1))
    payback_box.fill.solid()
    payback_box.fill.fore_color.rgb = PRIMARY_BLUE
    payback_box.line.fill.background()

    pb_text = slide.shapes.add_textbox(Inches(3.2), Inches(6.4), Inches(3.6), Inches(0.6))
    pb_frame = pb_text.text_frame
    pb_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    pbp1 = pb_frame.paragraphs[0]
    pbp1.text = "Payback Period: 2.1 Months"
    pbp1.font.size = Pt(20)
    pbp1.font.bold = True
    pbp1.font.color.rgb = WHITE
    pbp1.alignment = PP_ALIGN.CENTER

    # Slide 4: Implementation
    slide = add_content_slide(prs, "30-Day Implementation")

    weeks = [
        ("Week 1", "Discovery & Setup", "Kickoff, configuration, data prep", ACCENT_BLUE),
        ("Week 2", "Configuration", "System setup, workflows, testing", TEAL),
        ("Week 3", "Training", "Staff training, documentation", GREEN),
        ("Week 4", "Go-Live", "Launch, support, optimization", ORANGE)
    ]

    y_pos = Inches(1.5)
    for week, phase, activities, color in weeks:
        # Timeline marker
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y_pos, Inches(0.6), Inches(0.6))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

        m_text = slide.shapes.add_textbox(Inches(0.7), y_pos, Inches(0.6), Inches(0.6))
        m_frame = m_text.text_frame
        m_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        mp = m_frame.paragraphs[0]
        mp.text = week.split()[1]
        mp.font.size = Pt(20)
        mp.font.bold = True
        mp.font.color.rgb = WHITE
        mp.alignment = PP_ALIGN.CENTER

        # Content box
        content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(1.5), y_pos, Inches(8), Inches(1.2))
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = LIGHT_GRAY
        content_box.line.color.rgb = color
        content_box.line.width = Pt(2)

        # Phase name
        phase_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.15), Inches(7.6), Inches(0.4))
        p_frame = phase_box.text_frame
        pp = p_frame.paragraphs[0]
        pp.text = phase
        pp.font.size = Pt(18)
        pp.font.bold = True
        pp.font.color.rgb = color

        # Activities
        act_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.6), Inches(7.6), Inches(0.4))
        a_frame = act_box.text_frame
        ap = a_frame.paragraphs[0]
        ap.text = activities
        ap.font.size = Pt(13)
        ap.font.color.rgb = DARK_GRAY

        y_pos += Inches(1.5)

    prs.save('02-HOSPITAL-ADMINISTRATORS-EXECUTIVE-PRESENTATION.pptx')
    print("✓ Hospital Administrators presentation created!")

def create_medical_staff_presentation():
    """Create Medical Staff User Guide presentation"""
    print("Creating Medical Staff User Guide...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Your New Hospital Management System",
                   "Making Your Work Easier and More Efficient", TEAL)

    # Slide 2: Benefits Overview
    slide = add_content_slide(prs, "What's In It For You?", TEAL)

    roles_benefits = [
        ("👨‍⚕️ Doctors", "Save 1-2 hours/day\nInstant critical alerts\nComplete patient history", GREEN),
        ("👩‍⚕️ Nurses", "50% less phone time\nReal-time patient status\nBetter coordination", ACCENT_BLUE),
        ("📋 Reception", "Faster registration\nAuto queue placement\nLess paperwork", TEAL),
        ("🔬 Lab", "No phone tag\nClear work queue\nAuto notifications", ORANGE)
    ]

    x_start = Inches(0.5)
    for i, (role, benefits, color) in enumerate(roles_benefits):
        x = x_start + Inches(2.4) * i

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.5), Inches(2.2), Inches(3))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(1.7), Inches(1.8), Inches(2.6))
        t_frame = text_box.text_frame

        tp1 = t_frame.paragraphs[0]
        tp1.text = role
        tp1.font.size = Pt(16)
        tp1.font.bold = True
        tp1.font.color.rgb = WHITE
        tp1.alignment = PP_ALIGN.CENTER

        tp2 = t_frame.add_paragraph()
        tp2.text = "\n" + benefits
        tp2.font.size = Pt(11)
        tp2.font.color.rgb = LIGHT_GRAY
        tp2.alignment = PP_ALIGN.LEFT
        tp2.space_before = Pt(12)

    # Add key principle
    principle_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(1.5), Inches(5), Inches(7), Inches(1.5))
    principle_box.fill.solid()
    principle_box.fill.fore_color.rgb = LIGHT_GRAY
    principle_box.line.color.rgb = TEAL
    principle_box.line.width = Pt(3)

    principle_text = slide.shapes.add_textbox(Inches(1.7), Inches(5.2), Inches(6.6), Inches(1.1))
    pr_frame = principle_text.text_frame
    pr_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    prp1 = pr_frame.paragraphs[0]
    prp1.text = "One Patient, One Record, One System"
    prp1.font.size = Pt(28)
    prp1.font.bold = True
    prp1.font.color.rgb = TEAL
    prp1.alignment = PP_ALIGN.CENTER

    prp2 = pr_frame.add_paragraph()
    prp2.text = "Enter information once • Available everywhere • Real-time updates"
    prp2.font.size = Pt(14)
    prp2.font.color.rgb = DARK_GRAY
    prp2.alignment = PP_ALIGN.CENTER

    # Slide 3: Getting Started
    slide = add_content_slide(prs, "Getting Started", TEAL)

    steps = [
        ("1", "Login", "Use your employee email and password", ACCENT_BLUE),
        ("2", "Dashboard", "See your personalized work queue", TEAL),
        ("3", "Navigate", "Use the left sidebar menu", GREEN),
        ("4", "Notifications", "Click the bell icon for alerts", ORANGE),
        ("5", "Help", "Click '?' for context-specific help", PRIMARY_BLUE)
    ]

    y_pos = Inches(1.5)
    for num, step, desc, color in steps:
        # Number
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.5), Inches(0.5))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = color
        num_box.line.fill.background()

        n_text = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(0.5), Inches(0.5))
        n_frame = n_text.text_frame
        n_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = n_frame.paragraphs[0]
        np.text = num
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.alignment = PP_ALIGN.CENTER

        # Step box
        step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(1.5), y_pos, Inches(8), Inches(0.9))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = LIGHT_GRAY
        step_box.line.color.rgb = color
        step_box.line.width = Pt(2)

        # Content
        content_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.1), Inches(7.6), Inches(0.7))
        c_frame = content_box.text_frame

        cp1 = c_frame.paragraphs[0]
        cp1.text = step
        cp1.font.size = Pt(16)
        cp1.font.bold = True
        cp1.font.color.rgb = color

        cp2 = c_frame.add_paragraph()
        cp2.text = desc
        cp2.font.size = Pt(13)
        cp2.font.color.rgb = DARK_GRAY

        y_pos += Inches(1.05)

    prs.save('03-MEDICAL-STAFF-USER-GUIDE.pptx')
    print("✓ Medical Staff User Guide created!")

def create_technical_presentation():
    """Create Technical Stakeholders presentation"""
    print("Creating Technical Stakeholders presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Technical Architecture",
                   "Modern, Scalable, Secure Healthcare Platform", DARK_GRAY)

    # Slide 2: Technology Stack
    slide = add_content_slide(prs, "Technology Stack", DARK_GRAY)

    # Architecture layers
    layers = [
        ("Frontend", "React 18 • TypeScript • Material-UI\nResponsive • Real-time updates", ACCENT_BLUE),
        ("Backend", "Node.js 20 • Express • TypeScript\nRESTful APIs • Socket.IO", TEAL),
        ("Database", "PostgreSQL 14+ • ACID compliance\nRedis 7+ for caching", GREEN),
        ("Security", "JWT Auth • RBAC • Encryption\nHIPAA compliant", ORANGE)
    ]

    y_pos = Inches(1.5)
    for layer, tech, color in layers:
        # Layer box
        layer_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(1), y_pos, Inches(8), Inches(1.3))
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        layer_box.line.fill.background()

        # Layer name
        name_box = slide.shapes.add_textbox(Inches(1.3), y_pos + Inches(0.15), Inches(7.4), Inches(0.4))
        n_frame = name_box.text_frame
        np = n_frame.paragraphs[0]
        np.text = layer
        np.font.size = Pt(20)
        np.font.bold = True
        np.font.color.rgb = WHITE

        # Technologies
        tech_box = slide.shapes.add_textbox(Inches(1.3), y_pos + Inches(0.6), Inches(7.4), Inches(0.6))
        t_frame = tech_box.text_frame
        tp = t_frame.paragraphs[0]
        tp.text = tech
        tp.font.size = Pt(12)
        tp.font.color.rgb = LIGHT_GRAY

        y_pos += Inches(1.45)

    # Slide 3: Performance Metrics
    slide = add_content_slide(prs, "Performance & Scalability", DARK_GRAY)

    metrics = [
        ("1,000+", "Concurrent Users", GREEN),
        ("<200ms", "API Response Time", ACCENT_BLUE),
        ("99.9%", "Uptime SLA", TEAL),
        ("<50ms", "WebSocket Latency", ORANGE)
    ]

    x_positions = [Inches(0.5), Inches(2.8), Inches(5.1), Inches(7.4)]
    for i, (value, label, color) in enumerate(metrics):
        metric_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           x_positions[i], Inches(1.8), Inches(2.2), Inches(2))
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = LIGHT_GRAY
        metric_box.line.color.rgb = color
        metric_box.line.width = Pt(3)

        # Value
        val_box = slide.shapes.add_textbox(x_positions[i] + Inches(0.1), Inches(2.2), Inches(2), Inches(0.8))
        v_frame = val_box.text_frame
        v_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = v_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(36)
        vp.font.bold = True
        vp.font.color.rgb = color
        vp.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x_positions[i] + Inches(0.1), Inches(3.1), Inches(2), Inches(0.6))
        l_frame = lbl_box.text_frame
        l_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        lp = l_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(13)
        lp.font.color.rgb = DARK_GRAY
        lp.alignment = PP_ALIGN.CENTER

    # Key features
    features = [
        "✓ Horizontal scalability with stateless architecture",
        "✓ Redis caching for optimal performance",
        "✓ PostgreSQL connection pooling",
        "✓ Real-time WebSocket notifications",
        "✓ Auto-scaling based on load"
    ]

    y_pos = Inches(4.5)
    for feature in features:
        f_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(8), Inches(0.4))
        f_frame = f_box.text_frame
        fp = f_frame.paragraphs[0]
        fp.text = feature
        fp.font.size = Pt(15)
        fp.font.color.rgb = DARK_GRAY
        y_pos += Inches(0.45)

    prs.save('04-TECHNICAL-STAKEHOLDERS-ARCHITECTURE.pptx')
    print("✓ Technical Stakeholders presentation created!")

def create_patient_presentation():
    """Create Patient Experience presentation"""
    print("Creating Patient Experience presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Better Care, Faster Service",
                   "What Our New System Means For You", GREEN)

    # Slide 2: Benefits for Patients
    slide = add_content_slide(prs, "What's Changing?", GREEN)

    benefits = [
        ("⏱️", "40% Shorter Wait Times", "Less time in waiting room", GREEN),
        ("📋", "Less Paperwork", "Quick check-in process", TEAL),
        ("⚡", "Faster Results", "Test results to doctor in minutes", ACCENT_BLUE),
        ("🔒", "Your Privacy Protected", "Bank-level security", ORANGE)
    ]

    x_start = Inches(0.5)
    for i, (icon, benefit, desc, color) in enumerate(benefits):
        x = x_start + Inches(2.4) * i

        # Icon box
        icon_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), Inches(1.6), Inches(1), Inches(1))
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = color
        icon_box.line.fill.background()

        i_text = slide.shapes.add_textbox(x + Inches(0.6), Inches(1.6), Inches(1), Inches(1))
        i_frame = i_text.text_frame
        i_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        ip = i_frame.paragraphs[0]
        ip.text = icon
        ip.font.size = Pt(40)
        ip.alignment = PP_ALIGN.CENTER

        # Benefit text
        b_box = slide.shapes.add_textbox(x, Inches(2.8), Inches(2.2), Inches(0.6))
        b_frame = b_box.text_frame
        bp = b_frame.paragraphs[0]
        bp.text = benefit
        bp.font.size = Pt(14)
        bp.font.bold = True
        bp.font.color.rgb = color
        bp.alignment = PP_ALIGN.CENTER

        # Description
        d_box = slide.shapes.add_textbox(x, Inches(3.5), Inches(2.2), Inches(0.6))
        d_frame = d_box.text_frame
        dp = d_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(11)
        dp.font.color.rgb = DARK_GRAY
        dp.alignment = PP_ALIGN.CENTER

    # What stays the same
    same_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1), Inches(4.5), Inches(8), Inches(2))
    same_box.fill.solid()
    same_box.fill.fore_color.rgb = LIGHT_GRAY
    same_box.line.color.rgb = GREEN
    same_box.line.width = Pt(3)

    same_text = slide.shapes.add_textbox(Inches(1.3), Inches(4.8), Inches(7.4), Inches(1.4))
    s_frame = same_text.text_frame

    sp1 = s_frame.paragraphs[0]
    sp1.text = "What Stays the Same:"
    sp1.font.size = Pt(18)
    sp1.font.bold = True
    sp1.font.color.rgb = GREEN

    sp2 = s_frame.add_paragraph()
    sp2.text = "✓ Same caring doctors and nurses"
    sp2.font.size = Pt(14)
    sp2.font.color.rgb = DARK_GRAY

    sp3 = s_frame.add_paragraph()
    sp3.text = "✓ Same quality care you trust"
    sp3.font.size = Pt(14)
    sp3.font.color.rgb = DARK_GRAY

    sp4 = s_frame.add_paragraph()
    sp4.text = "✓ Your privacy remains completely protected"
    sp4.font.size = Pt(14)
    sp4.font.color.rgb = DARK_GRAY

    # Slide 3: Your Visit Experience
    slide = add_content_slide(prs, "Your Visit Experience - Before vs. After", GREEN)

    # Before column
    before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RED
    before_box.fill.fore_color.brightness = 0.7
    before_box.line.fill.background()

    before_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(4.1), Inches(0.5))
    bt_frame = before_title.text_frame
    btp = bt_frame.paragraphs[0]
    btp.text = "❌ Before"
    btp.font.size = Pt(24)
    btp.font.bold = True
    btp.font.color.rgb = RED
    btp.alignment = PP_ALIGN.CENTER

    before_text = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(3.7), Inches(3.8))
    btx_frame = before_text.text_frame

    items_before = [
        "Long paper forms to fill",
        "Wait 60-90 minutes",
        "No idea when your turn is",
        "Doctor searches for your chart",
        "Hand-written prescriptions",
        "Wait again at pharmacy",
        "Call in days for results"
    ]

    for item in items_before:
        p = btx_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # After column
    after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = GREEN
    after_box.fill.fore_color.brightness = 0.7
    after_box.line.fill.background()

    after_title = slide.shapes.add_textbox(Inches(5.7), Inches(1.7), Inches(3.6), Inches(0.5))
    at_frame = after_title.text_frame
    atp = at_frame.paragraphs[0]
    atp.text = "✓ Now"
    atp.font.size = Pt(24)
    atp.font.bold = True
    atp.font.color.rgb = GREEN
    atp.alignment = PP_ALIGN.CENTER

    after_text = slide.shapes.add_textbox(Inches(5.9), Inches(2.4), Inches(3.2), Inches(3.8))
    atx_frame = after_text.text_frame

    items_after = [
        "Quick 5-minute check-in",
        "Wait 30-45 minutes",
        "Know your queue position",
        "Complete history instantly",
        "Electronic prescriptions",
        "Meds ready when you arrive",
        "Same-day results often"
    ]

    for item in items_after:
        p = atx_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    prs.save('05-PATIENT-EXPERIENCE-OVERVIEW.pptx')
    print("✓ Patient Experience presentation created!")

if __name__ == "__main__":
    print("\n" + "="*60)
    print("Creating Beautiful Stakeholder Presentations")
    print("="*60 + "\n")

    create_hospital_admin_presentation()
    create_medical_staff_presentation()
    create_technical_presentation()
    create_patient_presentation()

    print("\n" + "="*60)
    print("✓ All presentations created successfully!")
    print("="*60 + "\n")
